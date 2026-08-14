"""多格式解析与 mime 分派单元测试。"""

from __future__ import annotations

import io

from app.modules.ingestion.parsers.dispatch import detect_format, parse_document_bytes
from app.modules.ingestion.parsers.docx_parser import parse_docx_bytes
from app.modules.ingestion.parsers.pptx_parser import parse_pptx_bytes
from app.modules.ingestion.parsers.text_parser import parse_text_bytes
from app.modules.ingestion.parsers.xlsx_parser import parse_xlsx_bytes


def test_parse_markdown_headings() -> None:
    data = b"# Title\n\n## Section\n\nhello world body text\n"
    pages = parse_text_bytes(data, mime_type="text/markdown")
    assert len(pages) == 1
    assert any(b.get("level") == 1 for b in pages[0].blocks)
    assert "hello world" in pages[0].plain_text


def test_parse_docx_paragraph_and_heading() -> None:
    from docx import Document

    buf = io.BytesIO()
    doc = Document()
    doc.add_heading("设备手册", level=1)
    doc.add_paragraph("液压站压力异常时先检查溢流阀。")
    doc.save(buf)
    pages = parse_docx_bytes(buf.getvalue())
    assert len(pages) == 1
    assert "液压站" in pages[0].plain_text
    assert any(b.get("level") == 1 for b in pages[0].blocks)


def test_parse_xlsx_sheet_as_page() -> None:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    assert ws is not None
    ws.title = "点检"
    ws.append(["设备", "周期"])
    ws.append(["HYD-2201", "月"])
    buf = io.BytesIO()
    wb.save(buf)
    pages = parse_xlsx_bytes(buf.getvalue())
    assert len(pages) == 1
    assert "HYD-2201" in pages[0].plain_text
    assert any(b.get("type") == "table" for b in pages[0].blocks)


def test_parse_pptx_slide() -> None:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    tf = box.text_frame
    tf.text = "开机前检查冷却水"
    for p in tf.paragraphs:
        for run in p.runs:
            run.font.size = Pt(18)
    buf = io.BytesIO()
    prs.save(buf)
    pages = parse_pptx_bytes(buf.getvalue())
    assert len(pages) == 1
    assert "冷却水" in pages[0].plain_text


def test_dispatch_routes_by_mime() -> None:
    md = b"# Hi\nbody"
    assert detect_format(md, "text/markdown") == "markdown"
    pages = parse_document_bytes(md, "text/markdown")
    assert pages[0].plain_text

    from docx import Document

    buf = io.BytesIO()
    Document().save(buf)
    data = buf.getvalue()
    assert (
        detect_format(
            data, "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )
        == "docx"
    )


def test_resolve_content_type_prefers_md_extension_over_text_plain() -> None:
    from app.modules.ingestion.parsers.dispatch import resolve_content_type

    assert resolve_content_type("text/plain", "手册.md") == "text/markdown"
    assert resolve_content_type("application/octet-stream", "a.markdown") == "text/markdown"
    assert resolve_content_type("text/plain", "original.md") == "text/markdown"
    assert resolve_content_type("application/pdf", "x.md") == "application/pdf"
