"""PPTX 解析（python-pptx）：每页幻灯片为一个逻辑页，含备注。"""

from __future__ import annotations

import io
from typing import Any

from app.modules.ingestion.parsers.common import page_from_blocks, text_block
from app.modules.ingestion.parsers.pdf import PageParse


def parse_pptx_bytes(data: bytes) -> list[PageParse]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(data))
    pages: list[PageParse] = []
    for idx, slide in enumerate(prs.slides, start=1):
        blocks: list[dict[str, Any]] = []
        order = 0
        for shape in slide.shapes:
            texts = _shape_texts(shape, MSO_SHAPE_TYPE)
            for text in texts:
                if not text.strip():
                    continue
                blocks.append(text_block(text, order=order))
                order += 1
        notes = _notes_text(slide)
        if notes.strip():
            blocks.append(text_block(f"备注：{notes}", order=order, size=11.0, bold=False))
        pages.append(page_from_blocks(idx, blocks))
    return pages or [page_from_blocks(1, [])]


def _shape_texts(shape: Any, shape_types: Any) -> list[str]:
    out: list[str] = []
    if shape.has_text_frame:
        for para in shape.text_frame.paragraphs:
            t = "".join(run.text for run in para.runs).strip() or (para.text or "").strip()
            if t:
                out.append(t)
    if shape.shape_type == shape_types.TABLE and shape.has_table:
        rows: list[list[str]] = []
        for row in shape.table.rows:
            rows.append([(" ".join(c.text.split())) for c in row.cells])
        if rows:
            header = rows[0]
            lines = [
                "| " + " | ".join(header) + " |",
                "| " + " | ".join("---" for _ in header) + " |",
            ]
            for r in rows[1:]:
                lines.append("| " + " | ".join(r) + " |")
            out.append("\n".join(lines))
    return out


def _notes_text(slide: Any) -> str:
    if not slide.has_notes_slide:
        return ""
    frame = slide.notes_slide.notes_text_frame
    if frame is None:
        return ""
    return (frame.text or "").strip()
